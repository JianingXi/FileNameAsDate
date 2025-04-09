import os
import re

# 允许处理的文件扩展名（全部转为小写判断）
ALLOWED_EXTENSIONS = {".txt", ".html", ".htm",
                      ".docx", ".doc",
                      ".pptx", ".ppt",
                      ".xlsx", ".xls"}

def read_text_file(file_path):
    """
    尝试以 UTF-8 编码读取文本文件，失败后使用 GBK 读取
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, 'r', encoding='gbk') as f:
                return f.read()
        except Exception as e:
            print(f"读取 {file_path} 时出错: {e}")
    except Exception as e:
        print(f"读取 {file_path} 时出错: {e}")
    return ""

def get_docx_text(file_path):
    """
    使用 python-docx 提取 docx 文件内容
    """
    try:
        import docx
    except ImportError:
        print("请安装 python-docx 库以解析 Word 文件: pip install python-docx")
        return ""
    try:
        doc = docx.Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n".join(paragraphs)
    except Exception as e:
        print(f"读取 {file_path} 时出错: {e}")
        return ""

def get_doc_text(file_path):
    """
    使用 win32com 提取 doc 文件内容（需要 Windows 环境和 Office 支持）
    """
    try:
        import win32com.client
    except ImportError:
        print("请安装 pywin32 库以解析 Word 文件: pip install pywin32")
        return ""
    text = ""
    word = None
    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(file_path, ReadOnly=True)
        text = doc.Content.Text
        doc.Close(False)
    except Exception as e:
        print(f"读取 {file_path} 时出错: {e}")
    finally:
        if word:
            word.Quit()
    return text

def get_pptx_text(file_path):
    """
    使用 python-pptx 读取 pptx 文件内容
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("请安装 python-pptx 库以解析 PPT 文件: pip install python-pptx")
        return ""
    try:
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        print(f"读取 {file_path} 时出错: {e}")
        return ""

def get_ppt_text(file_path):
    """
    使用 win32com 提取 ppt 文件内容（需要 Windows 环境和 Office 支持）
    """
    try:
        import win32com.client
    except ImportError:
        print("请安装 pywin32 库以解析 PPT 文件: pip install pywin32")
        return ""
    text = ""
    powerpoint = None
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = False
        presentation = powerpoint.Presentations.Open(file_path, WithWindow=False)
        texts = []
        for slide in presentation.Slides:
            for shape in slide.Shapes:
                if hasattr(shape, "HasTextFrame") and shape.HasTextFrame:
                    if shape.TextFrame.HasText:
                        texts.append(shape.TextFrame.TextRange.Text)
        text = "\n".join(texts)
        presentation.Close()
    except Exception as e:
        print(f"读取 {file_path} 时出错: {e}")
    finally:
        if powerpoint:
            powerpoint.Quit()
    return text

def get_xlsx_text(file_path):
    """
    使用 openpyxl 读取 xlsx 文件内容
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        print("请安装 openpyxl 库以解析 Excel 文件: pip install openpyxl")
        return ""
    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
        texts = []
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        texts.append(str(cell))
        wb.close()
        return "\n".join(texts)
    except Exception as e:
        print(f"读取 {file_path} 时出错: {e}")
        return ""

def get_xls_text(file_path):
    """
    使用 win32com 提取 xls 文件内容（需要 Windows 环境和 Office 支持）
    """
    try:
        import win32com.client
    except ImportError:
        print("请安装 pywin32 库以解析 Excel 文件: pip install pywin32")
        return ""
    text = ""
    excel = None
    try:
        excel = win32com.client.Dispatch("Excel.Application")
        excel.Visible = False
        workbook = excel.Workbooks.Open(file_path, ReadOnly=True)
        texts = []
        for sheet in workbook.Worksheets:
            used_range = sheet.UsedRange
            if used_range is None:
                continue
            for row in used_range.Rows:
                for cell in row.Cells:
                    if cell.Value is not None:
                        texts.append(str(cell.Value))
        text = "\n".join(texts)
        workbook.Close(False)
    except Exception as e:
        print(f"读取 {file_path} 时出错: {e}")
    finally:
        if excel:
            excel.Quit()
    return text

def get_file_text(file_path):
    """
    根据文件扩展名调用不同的解析方法获取文本内容
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext in {".txt", ".html", ".htm"}:
        return read_text_file(file_path)
    elif ext == ".docx":
        return get_docx_text(file_path)
    elif ext == ".doc":
        return get_doc_text(file_path)
    elif ext == ".pptx":
        return get_pptx_text(file_path)
    elif ext == ".ppt":
        return get_ppt_text(file_path)
    elif ext == ".xlsx":
        return get_xlsx_text(file_path)
    elif ext == ".xls":
        return get_xls_text(file_path)
    else:
        return ""

def extract_chinese_name(content, min_length=8):
    """
    从内容中匹配连续的中文字符序列，返回第一个长度不小于 min_length 的序列
    中文字符范围：[\\u4e00-\\u9fff]
    """
    pattern = re.compile(r'[\u4e00-\u9fff]+')
    matches = pattern.findall(content)
    for seq in matches:
        if len(seq) >= min_length:
            return seq[:min_length]
    return None

def get_unique_name(directory, new_name, ext):
    """
    在指定目录内判断新文件名是否冲突，如有冲突则在文件名末尾添加 _数字 后缀
    """
    candidate = new_name + ext
    counter = 1
    while os.path.exists(os.path.join(directory, candidate)):
        candidate = f"{new_name}_{counter}{ext}"
        counter += 1
    return candidate

def rename_files_by_content(root_dir):
    """
    递归扫描指定目录中所有文件，
    对于扩展名属于 ALLOWED_EXTENSIONS 且文件名（不含扩展名）长度小于 5 的文件，
    解析文件内容提取中文连续字符序列，
    如果提取到，则在原有文件名后追加"_"和提取的字符串作为新文件名（保留原扩展名）；
    若未提取到，则保持原有文件名不变。
    """
    for current_root, dirs, files in os.walk(root_dir):
        for filename in files:
            name, ext = os.path.splitext(filename)
            if ext.lower() not in ALLOWED_EXTENSIONS:
                continue
            if len(name) >= 6:
                continue
            file_path = os.path.join(current_root, filename)
            print(f"处理文件: {file_path}")
            content = get_file_text(file_path)
            append_str = extract_chinese_name(content, min_length=8)
            if append_str:
                # 新文件名为 原有文件名 + "_" + 提取的字符串
                new_base = f"{name}_{append_str}"
            else:
                # 如果没有提取到，则保持原有文件名
                new_base = name
            unique_new_filename = get_unique_name(current_root, new_base, ext)
            new_file_path = os.path.join(current_root, unique_new_filename)
            try:
                os.rename(file_path, new_file_path)
                print(f"已将 {file_path} 重命名为 {new_file_path}")
            except Exception as e:
                print(f"重命名 {file_path} 时出错: {e}")

if __name__ == "__main__":
    # 修改为你需要处理的目标目录路径
    target_directory = r"D:\Alpha\StoreLatestYears\Store2025\B教学_教学与人才培养_A03_学生竞赛"
    if not os.path.isdir(target_directory):
        print("错误: 指定的路径不是一个有效的目录!")
        exit(1)
    rename_files_by_content(target_directory)
    print("全部处理完成!")
